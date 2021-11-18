from pptx import Presentation
 
prs = Presentation()
 
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
 
title = slide.shapes.title
subtitle = slide.placeholders[1]
 
title.text = 'Hello, World!'
subtitle.text = 'python-pptx was here!'
 
prs.save('test.pptx')